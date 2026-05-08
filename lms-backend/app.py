from flask import Flask, request, jsonify
from flask_cors import CORS
from flask_sqlalchemy import SQLAlchemy
import os
import json
import re
from pptx import Presentation
from flask import send_file
import os
import nltk
import random
from flask import send_from_directory
nltk.download('punkt')
from transformers import pipeline

chatbot = pipeline(
    "text2text-generation",
    model="google/flan-t5-base"
)
app = Flask(__name__)
CORS(app)

# ================= DATABASE CONFIG =================
app.config['SQLALCHEMY_DATABASE_URI'] = 'sqlite:///lms.db'
app.config['SQLALCHEMY_TRACK_MODIFICATIONS'] = False

db = SQLAlchemy(app)


# ================= DATABASE MODELS =================

class User(db.Model):
    id = db.Column(db.Integer, primary_key=True)
    email = db.Column(db.String(100), unique=True)
    password = db.Column(db.String(100))
    role = db.Column(db.String(20))

class Quiz(db.Model):
    id = db.Column(db.Integer, primary_key=True)
    topic = db.Column(db.String(200))

class Question(db.Model):
    id = db.Column(db.Integer, primary_key=True)
    quiz_id = db.Column(db.Integer)
    question = db.Column(db.String(500))
    option1 = db.Column(db.String(200))
    option2 = db.Column(db.String(200))
    option3 = db.Column(db.String(200))
    option4 = db.Column(db.String(200))
    answer = db.Column(db.String(200))
class Material(db.Model):
    id = db.Column(db.Integer, primary_key=True)
    filename = db.Column(db.String(200))

class Video(db.Model):
    id = db.Column(db.Integer, primary_key=True)
    url = db.Column(db.String(300))
class QuizAttempt(db.Model):
    id = db.Column(db.Integer, primary_key=True)
    student_email = db.Column(db.String(100))
    quiz_id = db.Column(db.Integer)
    topic = db.Column(db.String(200))
    score = db.Column(db.Integer)
    total_questions = db.Column(db.Integer)

UPLOAD_FOLDER = "uploads"
os.makedirs(UPLOAD_FOLDER, exist_ok=True)
# ================= REGISTER =================

@app.route('/register', methods=['POST'])
def register():
    try:
        data = request.json

        existing = User.query.filter_by(email=data.get("email")).first()
        if existing:
            return jsonify({"message": "User already exists"})

        user = User(
            email=data.get("email"),
            password=data.get("password"),
            role=data.get("role")
        )

        db.session.add(user)
        db.session.commit()

        return jsonify({"message": "User registered successfully"})

    except Exception as e:
        print("REGISTER ERROR:", e)
        return jsonify({"message": "Error in registration"})


# ================= LOGIN =================

@app.route('/login', methods=['POST'])
def login():
    try:
        data = request.json

        user = User.query.filter_by(
            email=data.get("email"),
            password=data.get("password")
        ).first()

        if user:
            return jsonify({"role": user.role})

        return jsonify({"role": None})

    except Exception as e:
        print("LOGIN ERROR:", e)
        return jsonify({"role": None})


# ================= GENERATE QUIZ (AI) =================

import nltk
import random
from nltk.corpus import stopwords

nltk.download('punkt')
nltk.download('stopwords')

stop_words = set(stopwords.words('english'))

@app.route('/generate_quiz', methods=['POST'])
def generate_quiz():
    try:
        data = request.json
        text = data.get("text")

        if not text:
            return jsonify({"error": "No text provided"})

        # 🔥 LIMIT SIZE (prevents hanging)
        text = text[:2000]

        # 🔥 CLEAN TEXT (VERY IMPORTANT)
        text = re.sub(r"\*\*", "", text)
        text = re.sub(r"\(.*?\)", "", text)
        text = re.sub(r"[^a-zA-Z0-9\s\.]", " ", text)

        sentences = nltk.sent_tokenize(text)

        quiz = []
        used_keywords = set()

        for s in sentences:
            if len(quiz) == 10:
                break

            words = s.split()

            # ✅ Extract meaningful words
            candidates = []
            for w in words:
                clean = re.sub(r"[^\w]", "", w).lower()

                if clean and clean not in stop_words and len(clean) > 4:
                    candidates.append(clean)

            if len(candidates) < 4:
                continue

            keyword = random.choice(candidates)

            if keyword in used_keywords:
                continue

            used_keywords.add(keyword)

            # ✅ SAFE replacement
            question = s.replace(keyword, "______")

            # ✅ Create options
            options = list(set(candidates))
            random.shuffle(options)
            options = options[:4]

            if keyword not in options:
                options[0] = keyword

            random.shuffle(options)

            quiz.append({
                "question": question,
                "options": options,
                "answer": keyword
            })

        # 🔥 If not enough questions, repeat logic
        while len(quiz) < 10:
            quiz.append(random.choice(quiz))

        return jsonify({"questions": quiz})

    except Exception as e:
        print("QUIZ ERROR:", e)
        return jsonify({"error": str(e)})
# ================= SAVE QUIZ =================

@app.route('/save_quiz', methods=['POST'])
def save_quiz():
    try:
        data = request.json

        topic = data.get("topic")
        questions = data.get("questions")

        if not topic:
            return jsonify({"message": "Topic is required"}), 400

        if not questions:
            return jsonify({"message": "No questions to save"}), 400

        # ✅ Save quiz
        quiz = Quiz(topic=topic)
        db.session.add(quiz)
        db.session.commit()

        # ✅ Save questions
        for q in questions:
            question = Question(
                quiz_id=quiz.id,
                question=q["question"],
                option1=q["options"][0],
                option2=q["options"][1],
                option3=q["options"][2],
                option4=q["options"][3],
                answer=q["answer"]
            )
            db.session.add(question)

        db.session.commit()

        return jsonify({"message": "Quiz saved successfully"})

    except Exception as e:
        print("SAVE ERROR:", e)
        return jsonify({"message": str(e)})
# =================Generate PPT ============= #
@app.route('/generate_ppt', methods=['POST'])
def generate_ppt():
    try:
        text = request.json.get("text")

        if not text:
            return {"error": "No content provided"}

        # 🔥 CLEAN TEXT
        text = re.sub(r"\*\*", "", text)
        text = re.sub(r"\(.*?\)", "", text)

        sentences = text.split(".")
        sentences = [s.strip() for s in sentences if len(s.strip()) > 20]

        prs = Presentation()

        # 🎯 Title slide
        slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = "Presentation"
        slide.placeholders[1].text = "Auto Generated"

        # 🎯 FIXED STRUCTURE
        sections = [
            "Introduction",
            "Core Concept",
            "Types",
            "Applications",
            "Challenges",
            "Conclusion"
        ]

        index = 0

        for section in sections:
            if index >= len(sentences):
                break

            slide_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(slide_layout)

            slide.shapes.title.text = section

            tf = slide.placeholders[1].text_frame

            # ✅ Add max 3 short bullet points
            count = 0
            while count < 3 and index < len(sentences):
                sentence = sentences[index][:120]  # 🔥 limit length

                if count == 0:
                    tf.text = sentence
                else:
                    p = tf.add_paragraph()
                    p.text = sentence

                count += 1
                index += 1

        file_path = "generated_ppt.pptx"
        prs.save(file_path)

        return send_file(file_path, as_attachment=True)

    except Exception as e:
        print("PPT ERROR:", e)
        return {"error": str(e)}
# =============== upload material =========== #
@app.route('/upload_material', methods=['POST'])
def upload_material():
    try:
        file = request.files['file']

        filepath = os.path.join(UPLOAD_FOLDER, file.filename)
        file.save(filepath)

        material = Material(filename=file.filename)
        db.session.add(material)
        db.session.commit()

        return jsonify({"message": "Material uploaded successfully"})

    except Exception as e:
        print("UPLOAD ERROR:", e)
        return jsonify({"error": str(e)})
    
# ================upload video ==============#
@app.route('/upload_video', methods=['POST'])
def upload_video():
    try:
        data = request.json
        url = data.get("url")

        video = Video(url=url)
        db.session.add(video)
        db.session.commit()

        return jsonify({"message": "Video added successfully"})

    except Exception as e:
        print("VIDEO ERROR:", e)
        return jsonify({"error": str(e)})

# =========== student dashboard =========== #
@app.route('/search_materials')
def search_materials():
    query = request.args.get("query")

    if query:
        materials = Material.query.filter(
            Material.filename.ilike(f"%{query}%")
        ).all()
    else:
        materials = Material.query.all()

    return jsonify([
        {"filename": m.filename}
        for m in materials
    ])
# ============ search videos ======= #
@app.route('/search_videos')
def search_videos():
    query = request.args.get("query")

    if query:
        videos = Video.query.filter(
            Video.url.ilike(f"%{query}%")
        ).all()
    else:
        videos = Video.query.all()

    return jsonify([
        {"url": v.url} for v in videos
    ])
# =========== search quizzes =========== #
@app.route('/search_quizzes')
def search_quizzes():
    query = request.args.get("query")

    if query:
        quizzes = Quiz.query.filter(
            Quiz.topic.ilike(f"%{query}%")
        ).all()
    else:
        quizzes = Quiz.query.all()

    return jsonify([
        {
            "id": q.id,
            "topic": q.topic
        }
        for q in quizzes
    ])

@app.route('/uploads/<filename>')
def get_file(filename):
    return send_from_directory("uploads", filename)

# ============ attempt quiz ============= #
@app.route('/get_quiz/<int:quiz_id>')
def get_quiz(quiz_id):
    questions = Question.query.filter_by(quiz_id=quiz_id).all()

    return jsonify([
        {
            "id": q.id,
            "question": q.question,
            "option1": q.option1,
            "option2": q.option2,
            "option3": q.option3,
            "option4": q.option4
        }
        for q in questions
    ])
# =========== calculate score =========== #
@app.route('/submit_quiz', methods=['POST'])
def submit_quiz():
    data = request.json
    answers = data.get("answers")
    student_email = data.get("student_email")
    quiz_id = data.get("quiz_id")

    score = 0
    correct_answers = {}

    questions = Question.query.filter_by(quiz_id=quiz_id).all()

    for q in questions:
        selected = answers.get(str(q.id))

        correct_answers[q.id] = q.answer

        if selected == q.answer:
            score += 1

    quiz = Quiz.query.get(quiz_id)

    attempt = QuizAttempt(
        student_email=student_email,
        quiz_id=quiz_id,
        topic=quiz.topic,
        score=score,
        total_questions=len(questions)
    )

    db.session.add(attempt)
    db.session.commit()

    return jsonify({
        "score": score,
        "correct_answers": correct_answers
    })
# ============= student profile ============= #
@app.route('/student_profile')
def student_profile():
    email = request.args.get("email")
    print("PROFILE EMAIL:", email)
    attempts = QuizAttempt.query.filter_by(
        student_email=email
    ).all()
    print("FOUND:", len(attempts))
    total_correct = sum(a.score for a in attempts)
    total_questions = sum(a.total_questions for a in attempts)

    return jsonify({
        "username": email.split("@")[0],
        "tests_attempted": len(attempts),
        "total_correct": total_correct,
        "total_wrong": total_questions - total_correct,
        "quiz_history": [
            {
                "topic": a.topic,
                "score": a.score,
                "total": a.total_questions,
            }
            for a in attempts
        ]
    })
# ============ student performance for teacher =========== #
@app.route('/student_performance')
def student_performance():
    attempts = QuizAttempt.query.all()

    student_data = {}

    for attempt in attempts:
        email = attempt.student_email

        if email not in student_data:
            student_data[email] = {
                "tests_attempted": 0,
                "total_score": 0,
                "total_questions": 0
            }

        student_data[email]["tests_attempted"] += 1
        student_data[email]["total_score"] += attempt.score
        student_data[email]["total_questions"] += attempt.total_questions

    result = []

    for email, data in student_data.items():
        avg = (
            data["total_score"] / data["total_questions"]
        ) * 100 if data["total_questions"] > 0 else 0

        result.append({
            "email": email,
            "tests_attempted": data["tests_attempted"],
            "average_score": round(avg, 2),
            "status": (
                "Excellent" if avg >= 80
                else "Good" if avg >= 60
                else "Needs Improvement"
            )
        })

    return jsonify(result)

# =========== ai doubt solver ========= #
@app.route('/ask_doubt', methods=['POST'])
def ask_doubt():
    try:
        data = request.json
        question = data.get("question")

        result = chatbot(
            f"Answer this educational question clearly: {question}",
            max_length=150
        )

        return jsonify({
            "answer": result[0]["generated_text"]
        })

    except Exception as e:
        print(e)
        return jsonify({
            "answer": "Error generating answer"
        })
# ================= RUN APP =================

if __name__ == "__main__":
    app.run(debug=True)